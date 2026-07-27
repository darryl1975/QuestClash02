"""Plain-text extraction from uploaded source documents.

Each supported file extension maps to an extractor function that turns raw
bytes into plain text; extract_text() is the single entry point callers
should use (see main.py's /api/upload). All formats are read entirely into
memory — fine for the document sizes this app expects, but worth knowing if
very large uploads ever become a concern.
"""
import io
from pathlib import Path

import fitz  # PyMuPDF
import openpyxl
from docx import Document
from pptx import Presentation


def extract_text_from_txt(data: bytes) -> str:
    """Try common encodings in order; fall back to utf-8 with replacement
    chars so a bad encoding never hard-fails the upload."""
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def extract_text_from_pdf(data: bytes) -> str:
    doc = fitz.open(stream=data, filetype="pdf")
    return "\n".join(page.get_text() for page in doc)


def extract_text_from_docx(data: bytes) -> str:
    doc = Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_text_from_pptx(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
    return "\n".join(parts)


def extract_text_from_xlsx(data: bytes) -> str:
    """Flatten each sheet to '|'-joined rows, prefixed with a "Sheet: <name>"
    marker; data_only=True reads cached formula results rather than formulas."""
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts)


# File extension -> extractor function. .doc/.ppt share the .docx/.pptx
# handlers since python-docx/python-pptx can open the modern zip formats;
# legacy binary .doc/.ppt will fail here (not supported).
EXTRACTORS = {
    ".txt": extract_text_from_txt,
    ".pdf": extract_text_from_pdf,
    ".docx": extract_text_from_docx,
    ".doc": extract_text_from_docx,
    ".pptx": extract_text_from_pptx,
    ".ppt": extract_text_from_pptx,
    ".xlsx": extract_text_from_xlsx,
    ".xls": extract_text_from_xlsx,
}


def extract_text(filename: str, data: bytes) -> str:
    """Dispatch to the right extractor based on the filename's extension.
    Raises ValueError for anything not in EXTRACTORS."""
    ext = Path(filename).suffix.lower()
    extractor = EXTRACTORS.get(ext)
    if extractor is None:
        raise ValueError(f"Unsupported file type: {ext or '(none)'}")
    return extractor(data)
